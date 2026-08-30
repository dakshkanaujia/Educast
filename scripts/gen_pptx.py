#!/usr/bin/env python3
"""Fill the BITS Capstone 15-min demo template with EduCast content,
preserving the template's theme/branding."""
import copy
from pptx import Presentation

SRC = "Capstone_Project_15_Min_Demo_Template.pptx"
OUT = "docs/EduCast_Final_Presentation.pptx"

TEAM = ["Daksh Kanaujia (____________)",
        "K Manoj Krishna (____________)",
        "Utkersh Basnet (2023EBCS010)"]


def expand_names(shape, match, names):
    """Replace the paragraph equal to `match` with one paragraph per name,
    cloning the matched paragraph's XML so formatting is preserved."""
    tf = shape.text_frame
    for p in tf.paragraphs:
        if p.text.strip() == match:
            base = p._p
            anchor = base
            clones = [base]
            for _ in names[1:]:
                nc = copy.deepcopy(base)
                anchor.addnext(nc)
                anchor = nc
                clones.append(nc)
            for para_el, nm in zip(clones, names):
                runs = para_el.findall(
                    "{http://schemas.openxmlformats.org/drawingml/2006/main}r")
                if runs:
                    tnodes = runs[0].findall(
                        "{http://schemas.openxmlformats.org/drawingml/2006/main}t")
                    if tnodes:
                        tnodes[0].text = nm
                    for extra in runs[1:]:
                        para_el.remove(extra)
            return True
    return False

prs = Presentation(SRC)
slides = list(prs.slides)


def set_para_run_text(shape, match, new):
    """Replace the text of the paragraph whose text equals `match`,
    keeping the first run's formatting."""
    for p in shape.text_frame.paragraphs:
        if p.text.strip() == match:
            if p.runs:
                p.runs[0].text = new
                for r in p.runs[1:]:
                    r.text = ""
            else:
                p.text = new
            return True
    return False


def fill_body(shape, items):
    """Clear a content placeholder and refill with (text, level) bullets,
    inheriting theme fonts/sizes per outline level."""
    tf = shape.text_frame
    tf.clear()
    for i, (text, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = lvl


def body_shape(slide):
    # the content placeholder is the second shape with a text frame on 2..10
    cands = [s for s in slide.shapes if s.has_text_frame and s.name.startswith("Content")]
    return cands[0] if cands else None


def title_shape(slide):
    for s in slide.shapes:
        if s.has_text_frame and s.name.startswith("Title"):
            return s
    return None


# ---- Slide 1: Title ----
s1 = slides[0]
for sh in s1.shapes:
    if not sh.has_text_frame:
        continue
    set_para_run_text(sh, "Capstone Project Title",
                      "EduCast — A Demand-Driven, Real-Time Tutoring Marketplace")
    expand_names(sh, "Team Members", TEAM)
    set_para_run_text(sh, "Supervisor Name", "[Supervisor Name]")
    set_para_run_text(sh, "BSc CS (Academic Year)",
                      "BSc Computer Science (Online Mode) — 2025–2026")

# ---- Slide 2: Problem Statement ----
fill_body(body_shape(slides[1]), [
    ("Background of the problem", 0),
    ("Traditional e-learning (Udemy, Coursera, YouTube) is supply-driven: educators publish content speculatively and learners hope the right material exists.", 1),
    ("Gap in the existing system", 0),
    ("No fast way to request help for a specific, urgent problem — learners scroll catalogues never made for their exact question, at fixed prices with no competition.", 1),
    ("Importance of the problem", 0),
    ("Learners stuck on one specific thing need targeted, live help now; mentors want targeted, well-paid work matched to their expertise.", 1),
])

# ---- Slide 3: Objectives & Scope ----
fill_body(body_shape(slides[2]), [
    ("Objectives:", 0),
    ("Build a demand-driven marketplace: students post bounties, mentors bid", 1),
    ("Real-time events over WebSocket; price negotiation via counter-offers", 1),
    ("Atomic bid acceptance (no double-booking) + simulated escrow → release", 1),
    ("Persisted ratings/reviews, mentor profiles; one-command Docker deployment", 1),
    ("Scope:", 0),
    ("In-scope: auth + RBAC, bounty/bid lifecycle, negotiation, session + chat, completion + rating, mentor directory/profile, price insight", 1),
    ("Out-of-scope: real payments, live video/audio, push notifications", 1),
])

# ---- Slide 4: Existing System / Literature Review ----
fill_body(body_shape(slides[3]), [
    ("Existing approach 1: MOOC platforms (Udemy, Coursera) — pre-recorded, supply-driven catalogues", 0),
    ("Existing approach 2: Q&A / tutoring (Chegg, Wyzant, YouTube) — fixed pricing or async answers", 0),
    ("Limitations", 0),
    ("No demand-driven requests; no reverse-auction price discovery", 1),
    ("No built-in negotiation; weak real-time matching for urgent, specific needs", 1),
])

# ---- Slide 5: Proposed System Architecture ----
fill_body(body_shape(slides[4]), [
    ("System overview", 0),
    ("Three-tier: React Native (Expo) client ↔ Go/Gin REST + WebSocket hub ↔ PostgreSQL, secured by JWT auth + role-based access", 1),
    ("Architecture diagram", 0),
    ("[Insert architecture diagram — see Project Report §2.1]", 1),
    ("Module description", 0),
    ("Auth, Bounty, Bid, Negotiation, Acceptance/Session, Completion/Review, Mentor directory, and a central real-time Hub (targeted + broadcast messaging)", 1),
])

# ---- Slide 6: Tools & Technologies ----
fill_body(body_shape(slides[5]), [
    ("Programming Language: Go (backend); JavaScript / React Native (frontend)", 0),
    ("Frameworks: Gin, GORM, React Native (Expo), React Navigation", 0),
    ("Database: PostgreSQL 16", 0),
    ("Tools: Gorilla WebSocket, JWT, bcrypt, Docker & Docker Compose, Git/GitHub", 0),
])

# ---- Slide 7: Implementation / Demo ----
fill_body(body_shape(slides[6]), [
    ("Feature 1: Real-time bounty → bid → accept loop pushed live over WebSocket", 0),
    ("Feature 2: Price negotiation (counter-offers), atomic acceptance, simulated escrow/release, ratings & mentor profiles", 0),
    ("Screenshots / Flow", 0),
    ("[Insert app screenshots + demo video link]", 1),
])

# ---- Slide 8: Results & Analysis ----
fill_body(body_shape(slides[7]), [
    ("Output: full marketplace loop working end-to-end on the Dockerised stack", 0),
    ("Performance / validation: 33/33 automated API test cases pass (100%); API responses in sub-ms–low-ms (Gin logs); real-time events delivered instantly", 0),
    ("Comparison: vs supply-driven catalogues — demand-first, negotiated, real-time, with built-in accountability (atomic accept, escrow, reputation)", 0),
])

# ---- Slide 9: Challenges & Limitations ----
fill_body(body_shape(slides[8]), [
    ("Technical challenges", 0),
    ("Datastore migration MySQL → PostgreSQL (driver, DSN, schema rewrite)", 1),
    ("Docker orchestration with DB healthcheck gating the backend", 1),
    ("WebSocket hub design: targeted vs broadcast delivery; live presence", 1),
    ("Limitations", 0),
    ("Simulated payments; no live video; ephemeral chat; events reach only connected clients", 1),
])

# ---- Slide 10: Conclusion & Future Work ----
fill_body(body_shape(slides[9]), [
    ("Conclusion:", 0),
    ("A working demand-driven, real-time tutoring marketplace; UVP validated; 33/33 tests pass; fully containerised", 1),
    ("Future Work:", 0),
    ("Real payments + wallet; WebRTC live sessions; persistent notifications + chat history", 1),
    ("Bounty upvoting; search/pagination; automated Go tests in CI", 1),
])

prs.save(OUT)
print("Saved", OUT, "with", len(slides), "slides")
